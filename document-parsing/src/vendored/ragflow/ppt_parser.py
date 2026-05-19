"""Vendored from ragflow/deepdoc/parser/ppt_parser.py — verbatim (no rag.* deps)."""
from __future__ import annotations

import logging
from io import BytesIO

from pptx import Presentation


class RAGFlowPptParser:
    def __init__(self):
        super().__init__()
        self._shape_cache = {}

    def __sort_shapes(self, shapes):
        cache_key = id(shapes)
        if cache_key not in self._shape_cache:
            self._shape_cache[cache_key] = sorted(
                shapes,
                key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0),
            )
        return self._shape_cache[cache_key]

    def __get_bulleted_text(self, paragraph):
        is_bulleted = (
            bool(paragraph._p.xpath("./a:pPr/a:buChar"))
            or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum"))
            or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        )
        if is_bulleted:
            return f"{'  ' * paragraph.level}.{paragraph.text}"
        return paragraph.text

    def __extract(self, shape):
        try:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                texts = []
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "\n".join(texts)
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                if hasattr(shape, "text"):
                    return shape.text.strip()
                return ""
            if shape_type == 19:  # table
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    rows.append(
                        "; ".join(
                            [
                                tb.cell(0, j).text + ": " + tb.cell(i, j).text
                                for j in range(len(tb.columns))
                                if tb.cell(i, j)
                            ]
                        )
                    )
                return "\n".join(rows)
            if shape_type == 6:  # group
                texts = []
                for p in self.__sort_shapes(shape.shapes):
                    t = self.__extract(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)
            return ""
        except Exception as e:
            logging.error("Error processing shape: %s", e)
            return ""

    def __call__(self, fnm, from_page, to_page, callback=None):
        ppt = Presentation(fnm) if isinstance(fnm, str) else Presentation(BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in self.__sort_shapes(slide.shapes):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))
        return txts
