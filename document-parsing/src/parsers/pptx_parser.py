"""PPTX → markdown.

Walks slides directly to render structure faithfully:
  * slide title → ``## Slide N: <title>``
  * bullets → ``- ...``
  * tables → markdown tables
  * pictures → ``![name](data-uri)`` (extracted as images on the side)
"""
from __future__ import annotations

import base64
import logging
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.base import BaseParser, ExtractedImage, ParseResult

log = logging.getLogger(__name__)


class PptxParser(BaseParser):
    name = "pptx"
    extensions = ("pptx",)

    def parse(self, payload: bytes, filename: str) -> ParseResult:
        prs = Presentation(BytesIO(payload))
        out: list[str] = []
        images: list[ExtractedImage] = []
        image_idx = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            shapes = sorted(
                slide.shapes,
                key=lambda s: (
                    (s.top if s.top is not None else 0) // 10,
                    s.left if s.left is not None else 0,
                ),
            )
            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip()

            out.append(f"## Slide {slide_idx}" + (f": {title_text}" if title_text else ""))

            for shape in shapes:
                if shape == slide.shapes.title:
                    continue
                rendered = self._render_shape(shape)
                if rendered:
                    out.append(rendered)

                # Picture extraction (separate from text rendering).
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        image_idx += 1
                        suffix = (img.ext or "png").lower()
                        name = f"slide{slide_idx}_img{image_idx}.{suffix}"
                        mime = f"image/{'jpeg' if suffix in ('jpg', 'jpeg') else suffix}"
                        images.append(ExtractedImage(
                            name=name,
                            bytes_b64=base64.b64encode(img.blob).decode("ascii"),
                            mime=mime,
                        ))
                        out.append(f"![{name}]({name})")
                    except Exception as e:
                        log.debug("skip picture on slide %d: %s", slide_idx, e)

        return ParseResult(
            markdown="\n\n".join(out).strip(),
            parser=self.name,
            page_count=len(prs.slides),
            metadata={"extension": "pptx", "slide_count": len(prs.slides)},
            images=images,
        )

    def _render_shape(self, shape) -> str:
        try:
            if shape.has_text_frame:
                lines: list[str] = []
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if not txt:
                        continue
                    is_bulleted = (
                        bool(para._p.xpath("./a:pPr/a:buChar"))
                        or bool(para._p.xpath("./a:pPr/a:buAutoNum"))
                        or bool(para._p.xpath("./a:pPr/a:buBlip"))
                    )
                    if is_bulleted:
                        indent = "  " * para.level
                        lines.append(f"{indent}- {txt}")
                    else:
                        lines.append(txt)
                return "\n".join(lines)
        except Exception:
            pass

        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return _table_to_md(shape.table)
        except Exception:
            pass

        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                parts = [self._render_shape(s) for s in shape.shapes]
                return "\n\n".join(p for p in parts if p)
        except Exception:
            pass

        return ""


def _table_to_md(table) -> str:
    rows = []
    for row in table.rows:
        rows.append([cell.text.replace("\n", " ").strip() for cell in row.cells])
    rows = [r for r in rows if any(c for c in r)]
    if not rows:
        return ""
    n_cols = max(len(r) for r in rows)
    rows = [r + [""] * (n_cols - len(r)) for r in rows]
    header, body = rows[0], rows[1:]
    lines = ["| " + " | ".join(c.replace("|", r"\|") for c in header) + " |",
             "| " + " | ".join(["---"] * n_cols) + " |"]
    for r in body:
        lines.append("| " + " | ".join(c.replace("|", r"\|") for c in r) + " |")
    return "\n".join(lines)
